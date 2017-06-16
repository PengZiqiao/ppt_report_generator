from datetime import date, timedelta

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

import pandas as pd


def analyze_ppt(input, output):
    prs = Presentation(input)

    for index, layout in enumerate(prs.slide_layouts):
        slide = prs.slides.add_slide(layout)
        try:
            title = slide.shapes.title
            title.text = f'page-{index}:title'
        except AttributeError:
            print(f'>>> page-{index}:no title')

        field = slide.placeholders
        for each in field:
            each.text = str(each.placeholder_format.idx)

    prs.save(output)


def df_to_table(shape, df):
    # 插入表格
    rows, cols = df.shape
    tb = shape.insert_table(rows + 1, cols).table

    # 填写表头
    colnames = list(df.columns)
    for col_index, col_name in enumerate(colnames):
        tb.cell(0, col_index).text = col_name

    # 填写数据
    m = df.as_matrix()
    for row in range(rows):
        for col in range(cols):
            val = m[row, col]
            text = str(val)
            tb.cell(row + 1, col).text = text


def chart(df):
    data = ChartData()
    data.categories = list(df.index)
    for col in df:
        data.add_series(col, df[col].tolist())
    type = XL_CHART_TYPE.COLUMN_CLUSTERED
    return type, data


class Week:
    sunday = date.today()
    while not sunday.weekday() == 0:
        sunday -= timedelta(days=1)
    monday = sunday - timedelta(days=6)
    year = monday.year
    date_range = f"{monday.month}月{monday.day}日-{sunday.month}月{sunday.day}日"
    nw = f"第{monday.strftime('%U')}周"


class Report:
    def __init__(self, inputfile):
        """
        :param inputfile: file path of the template.pptx
        """
        self.prs = Presentation(inputfile)

    def liangjia(self, wuye, df1, df2):
        # 使用模板文件中母版第一页创建一个page
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        field = slide.placeholders
        # 标题
        field[0].text = f'{wuye}市场-量价'
        # 结论
        field[10].text = f'本周{wuye}市场整体表现为……'
        # 左图标题
        field[11].text = '南京（不含高淳溧水）近10周商品住宅市场供销量价'
        # 左图
        chart_type, chart_data = chart(df1)
        field[12].insert_chart(chart_type, chart_data)
        # 右图标题
        week = Week()
        field[13].text = f'{week.year}年{week.nw}南京（不含高淳溧水）商品住宅市场分板块供销量价'
        # 右图
        chart_type, chart_data = chart(df2)
        field[14].insert_chart(chart_type, chart_data)
        # 说理
        field[15].text = f'结论1\r结论2\r结论3'
        return self

    def paihang(self, wuye, df1, df2):
        # 使用模板文件中母版第二页创建一个page
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        field = slide.placeholders
        # 标题
        field[0].text = f'{wuye}市场-排行榜'
        # 结论
        field[10].text = f'xx项目销售xx套住宅，共xx万㎡……'
        # 左表标题
        week = Week()
        year = week.year
        nw = week.nw
        date_range = week.date_range
        field[11].text = f'{year}年{nw}（{date_range}）销售面积排行榜'
        # 左表
        df_to_table(field[14], df1)
        # 右图标题
        field[13].text = f'{year}年{nw}（{date_range}）销售套数排行榜'
        # 右表
        df_to_table(field[15], df2)
        return self

    def save(self, outputfile):
        self.prs.save(outputfile)


if __name__ == '__main__':
    analyze_ppt('周报母版.pptx', 'out.pptx')
    r = Report('周报母版.pptx')
    df1 = pd.read_excel('dataframes.xlsx', '00df1', index_col=0)
    df2 = pd.read_excel('dataframes.xlsx', '00df2', index_col=0)
    r.liangjia('住宅', df1, df2)
    df1 = pd.read_excel('dataframes.xlsx', '01df1')
    df2 = pd.read_excel('dataframes.xlsx', '01df2')
    r.paihang('住宅', df1, df2)
    # r.liangjia('商业')
    # r.liangjia('办公')
    # r.liangjia('别墅')
    r.save('output.pptx')
