from pptx import Presentation

path = 'input.pptx'
pres = Presentation(path)
i = 0
for slide in pres.slides:
    for shape in slide.shapes:
        try:
            shape.text = str(i)
            i += 1
        except AttributeError:
            print('cant change text')

pres.save('output.pptx')